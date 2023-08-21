import openai
import os
from flask import Flask, send_file
from flask_sqlalchemy import SQLAlchemy
from flask_login import LoginManager
from flask_migrate import Migrate
from functools import wraps
import base64
import psycopg2
from io import BytesIO
from psycopg2 import extras
from datetime import date
import json
import PyPDF2
from pptx import Presentation



#imports of routes.py
from flask import request, render_template, flash, redirect, url_for, flash, abort, session, jsonify, Response, render_template_string
from werkzeug.urls import url_parse
from flask_login import current_user, login_user, logout_user, login_required

#imports of forms.py
from flask_wtf import FlaskForm
from flask_wtf.file import FileField, FileAllowed, FileRequired
from wtforms import StringField, PasswordField, BooleanField, SubmitField, RadioField, DateField, SelectField, TextAreaField, TimeField
from wtforms.validators import ValidationError, DataRequired, Email, EqualTo, Regexp

#imports from models.py
from datetime import datetime
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
from flask_login import UserMixin



#instantiate application and database
app = Flask(__name__)
UPLOAD_FOLDER = 'static/pdfFiles/'
app.config['SECRET_KEY'] = 'edem-and-antonio'
app.config['SQLALCHEMY_DATABASE_URI'] = 'postgresql://postgres:password@localhost/lms_chatbot'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 30 * 1024 * 1024  # 30 MB limit
openai.api_key = os.environ.get('OPEN_AI_KEY')

db = SQLAlchemy(app)
migrate = Migrate(app, db)

#create login manager
login_manager = LoginManager()
login_manager.init_app(app)




ALLOWED_EXTENSIONS = {'pdf', 'pptx', 'ppt'}

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def get_flashed_messages(with_categories=False):
    """
    Custom implementation of get_flashed_messages to support categories.
    """
    flashes = app.jinja_env.globals['flashes']
    app.jinja_env.globals['flashes'] = []
    return flashes if with_categories else [msg for msg, _ in flashes]

def student_required(func):
    @wraps(func)
    def decorated_view(*args, **kwargs):
        if not current_user.is_authenticated or current_user.role != 'student':
            abort(403)  # Return a forbidden error if the user is not authenticated as a user
        return func(*args, **kwargs)
    return decorated_view


def lecturer_required(func):
    @wraps(func)
    def decorated_view(*args, **kwargs):
        if not current_user.is_authenticated or current_user.role != 'lecturer':
            abort(403)  # Return a forbidden error if the user is not authenticated as an organizer
        return func(*args, **kwargs)
    return decorated_view


def admin_required(func):
    @wraps(func)
    def decorated_view(*args, **kwargs):
        if not current_user.is_authenticated or current_user.role != 'admin':
            abort(403)  # Return a forbidden error if the user is not authenticated as an organizer
        return func(*args, **kwargs)
    return decorated_view

def create_quiz_objects(data, course_name):
    quiz = Quiz(title=course_name)  # You can adjust the title as needed
    db.session.add(quiz)
    db.session.commit()
    for item in data:
        question = Question(text=item['question'], quiz=quiz)
        db.session.add(question)
        db.session.commit()
        for answer_text in item['answers']:
            is_correct = answer_text == item['correct_answer']
            answer = Answer(text=answer_text, is_correct=is_correct, question=question)
            db.session.add(answer)
            db.session.commit()









class Lecturer(UserMixin, db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(100), nullable=False)
    surname  = db.Column(db.String(100), nullable=False)
    email = db.Column(db.String(100), nullable=False)
    phone_number = db.Column(db.String(100), nullable=False)
    password_hash = db.Column(db.String(1000), nullable=False)
    ghana_card = db.Column(db.String(20), nullable=False)
    role = db.Column(db.String(100), nullable=False)
    courses = db.relationship('Course', backref='lecturer', lazy=True)
    
    
    def __repr__(self):
        return '<User {}>'.format(self.username)
    def set_password(self, password):
        self.password_hash = generate_password_hash(password)
    def check_password(self, password):
        return check_password_hash(self.password_hash, password)

class Student(UserMixin, db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(100), nullable=False)
    surname = db.Column(db.String(100), nullable=False)
    email  = db.Column(db.String(100), nullable=False)
    phone_number = db.Column(db.String(100), nullable=False)
    password_hash = db.Column(db.String(1000), nullable=False)
    ghana_card = db.Column(db.String(20), nullable=False)
    avatar = db.Column(db.LargeBinary, nullable=True)
    role = db.Column(db.String(100), nullable=False)
    enrollments = db.relationship('Enrollment', backref='student', lazy=True)
    
    def __repr__(self):
        return '<User {}>'.format(self.username)
    def set_password(self, password):
        self.password_hash = generate_password_hash(password)
    def check_password(self, password):
        return check_password_hash(self.password_hash, password)

class Admin(UserMixin, db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(100), nullable=False)
    password_hash = db.Column(db.String(1000), nullable=False)
    role = db.Column(db.String(100), nullable=False)
    
    def __repr__(self):
        return '<User {}>'.format(self.username)
    def set_password(self, password):
        self.password_hash = generate_password_hash(password)
    def check_password(self, password):
        return check_password_hash(self.password_hash, password)



class Enrollment(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    student_id = db.Column(db.Integer, db.ForeignKey('student.id'), nullable=False)
    course_id = db.Column(db.Integer, db.ForeignKey('course.id'), nullable=False)

class Module(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(100), nullable=False)
    course_id = db.Column(db.Integer, db.ForeignKey('course.id'),nullable=False)
    lessons = db.relationship('Lesson', backref='module', lazy=True)


class Lesson(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    title = db.Column(db.String(100), nullable=False)
    module_id = db.Column(db.Integer, db.ForeignKey('module.id'), nullable=False)

class Course(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(100), nullable=False)
    lecturer_id = db.Column(db.Integer, db.ForeignKey('lecturer.id'), nullable=False,)
    modules = db.relationship('Module', backref='course', lazy=True)

class Section(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(100), nullable=True)

    

      



class StudentCourseProgress(db.Model):
    progress_id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('student.id'))
    course_id = db.Column(db.Integer, db.ForeignKey('course.id'))
    score = db.Column(db.Integer)
    completion_status = db.Column(db.Boolean)

class StudentModuleProgress(db.Model):
    progress_id = db.Column(db.Integer, primary_key=True)
    student_id = db.Column(db.Integer, db.ForeignKey('student.id'))
    course_id = db.Column(db.Integer, db.ForeignKey('course.id'))
    score =  db.Column(db.Integer)
    completion_status = db.Column(db.Boolean)

class StudentLessonProgress(db.Model):
    progress_id = db.Column(db.Integer, primary_key=True)
    student_id = db.Column(db.Integer, db.ForeignKey('student.id'))
    course_id = db.Column(db.Integer, db.ForeignKey('course.id'))
    score = db.Column(db.Integer)
    completion_status = db.Column(db.Boolean)

class StudentSectionProgress(db.Model):
    progress_id = db.Column(db.Integer, primary_key=True)
    student_id = db.Column(db.Integer, db.ForeignKey('student.id'))
    course_id = db.Column(db.Integer, db.ForeignKey('course.id'))
    score = db.Column(db.Integer)
    completion_status = db.Column(db.Boolean)


    




class Quiz(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    title = db.Column(db.String(100), nullable=False)
    questions = db.relationship('Question', backref='quiz', lazy=True)


class Question(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    text = db.Column(db.String(255), nullable=False)
    quiz_id = db.Column(db.Integer, db.ForeignKey('quiz.id'), nullable=False)
    answers = db.relationship('Answer', backref='question', lazy=True)


class Answer(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    text = db.Column(db.String(255), nullable=False)
    is_correct = db.Column(db.Boolean, default=False)
    question_id = db.Column(db.Integer, db.ForeignKey('question.id'), nullable=False)

class Submission(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    student_id = db.Column(db.Integer, db.ForeignKey('student.id'), nullable=False)
    quiz_id = db.Column(db.Integer, db.ForeignKey('quiz.id'), nullable=False)
    submission_details = db.relationship('SubmissionDetail', backref='submission', lazy=True)

class Grade(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    submission_id = db.Column(db.Integer, db.ForeignKey('submission.id'), nullable=False)
    grade = db.Column(db.Float, nullable=False)



class SubmissionDetail(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    submission_id = db.Column(db.Integer, db.ForeignKey('submission.id'), nullable=False)
    question_id = db.Column(db.Integer, db.ForeignKey('question.id'), nullable=False)
    selected_answer_id = db.Column(db.Integer, db.ForeignKey('answer.id'))
    is_correct = db.Column(db.Boolean, default=False)


def extract_text_from_pdf(pdf_file_path):
    text = ""
    with open(pdf_file_path, 'rb') as pdf_file:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text += page.extract_text()
    return text



def extract_text_from_pptx(pptx_file):
    presentation = Presentation(pptx_file)
    extracted_text = ""

    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text += shape.text + "\n"

    return extracted_text

def generate_text(message):
    completion = openai.ChatCompletion.create(
    model="gpt-3.5-turbo",
    messages=[
        {"role": "user", "content": message}
    ]
    )
    if completion.choices[0].message!=None:
        return completion.choices[0].message.content

    else :
        return 'Failed to Generate response!'



@login_manager.user_loader
def load_user(user_id):
    if session.get('my_role') == 'student_role':
        return  Student.query.get(int(user_id))
    elif session.get('my_role') == 'admin_role':
        return Admin.query.get(int(user_id))
    return Lecturer.query.get(int(user_id))


# ROUTES
@app.route('/', methods=['POST', 'GET'])
def index():
    template_string = '''
    <form action="/" method="POST">
    <input type="text" name="promptings" id="promtings" />
    <input type="submit" value="submit"/>
    </form>
    <div>{{ response }} {{ message }}</div>
    '''
    if request.method == "POST":
        prompt  = request.form.get('promptings')
        prompt_results = generate_text(prompt)
        return render_template_string(template_string, response=prompt_results)
    else:
        try:
            print(prompt_results)
            if prompt_results == "":
                raise NameError
        except:
            message = "keep beeping"
            prompt_results = ""
        return render_template('index.html', response="", message=message)


    
@app.route('/login', methods=["POST", "GET"])
def login():
    if request.method == "GET":
        return render_template('login.html')
    else:
        email  = request.form.get('email')
        password= request.form.get('password')
        account_type = request.form.get('user_type')
        if account_type == 'student':
            session['my_role'] = 'student_role'
            student = Student.query.filter_by(email=email).first()
            if student and student.check_password(password):
                login_user(student, remember=False)
                next_page = request.args.get('next')
                return redirect(next_page) if next_page else redirect(url_for('student_profile', _external=True))
            flash("You entered incorrect data", category='error')
            return redirect('login')
        else:
            session['my_role'] = 'lecture_role'
            lecturer = Lecturer.query.filter_by(email=email).first()
            if lecturer and lecturer.check_password(password):
                login_user(lecturer, remember=False)
                next_page = request.args.get('next')
                return redirect(next_page) if next_page else redirect(url_for('lecturer_profile', _external=True))

@app.route('/register', methods=['POST', 'GET'])
def register():
    if request.method == "GET":
        return render_template("register.html")
    else:
        first_name = request.form.get('first_name')
        surname = request.form.get('surname')
        account_type = request.form.get('user_type')
        email  = request.form.get('email')
        phone_number = request.form.get('phone_number')
        ghana_card = request.form.get('ghana_card')
        password_one = request.form.get("password")
        password_two = request.form.get("password_confirmation")
        if password_one != password_two:
            flash('Password entered is incorrect', category='error')
            return redirect('register')
        if account_type == "student":
            new_student = Student(name=first_name, surname=surname, email=email, phone_number=phone_number, ghana_card=ghana_card,role="student")
            new_student.set_password(password_one)
            db.session.add(new_student)
            db.session.commit()
            flash('Student Account has been created sucessfully', category='success')
            return redirect('login')
        else:
            new_lecturer = Lecturer(name=first_name, surname=surname, email=email, phone_number=phone_number, ghana_card=ghana_card, role="lecturer")
            new_lecturer.set_password(password_one)
            db.session.add(new_lecturer)
            db.session.commit()
            flash('Lecturer Account has been created successfully', category='success')
            return redirect('login')




@app.route('/lecturer_profile')
@login_required
@lecturer_required
def lecturer_profile():
    return render_template('dashboard.html')


@app.route('/student_profile')
@login_required
@student_required
def student_profile():
    quizzes = Quiz.query.all()
    print(quizzes)
    print(quizzes[-1].questions)
    print(quizzes[-1].questions[0].text)
    print(quizzes[-1].questions[0].answers)
    print(quizzes[-1].questions[0].answers[0])
    print(quizzes[-1].questions[0].answers[0].text)
    return render_template('courses.html', quizzes=quizzes)


@app.route('/start_quiz/<int:quiz_id>')
def start_quiz(quiz_id):
    quiz = Quiz.query.filter_by(id=quiz_id).first()
    questions = quiz.questions
    print(quiz.title)
    return render_template('quiz.html', quiz=quiz, questions=questions)


@app.route('/create_course_pdf', methods=['POST', 'GET'])
def create_course_pdf():
    if request.method == "POST":
        file = request.files['slide']
        file_extension = file.filename.rsplit('.', 1)[-1] if '.' in file.filename else None
        print(file_extension)
        if not allowed_file(file.filename):
            return jsonify({'error': 'Invalid file extension. Only PDF and PPTX files are allowed.'}), 400
        
        filename = secure_filename(file.filename)
        print(os.path.join(app.config['UPLOAD_FOLDER'], filename))
        file.save(os.path.join(app.config['UPLOAD_FOLDER'], filename))
        slide_file = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        if file_extension == "pdf":
            extracted_text = extract_text_from_pdf(slide_file)
        else:
            extracted_text = extract_text_from_pptx(slide_file)
        
        # print(extracted_text)
        num_quiz = request.form.get('quizzes')
        prompt = """Analyze the quoted text ""%s"" and generate %s questions  as a json object with the keys question, answers, correct_answer 
it should be in {} enclosed in a python list. I WANT ONLY THE OUTPUT""" % (extracted_text, num_quiz)
        prompt_results = generate_text(prompt)
        # quiz_dict =  json.loads(prompt_results)
        # for x in quiz_dict:
        #     print(x['question'])
        #     print(x.keys())
        quiz_dict = eval(prompt_results)
        create_quiz_objects(quiz_dict, filename)
        
        print(type(quiz_dict))
        return quiz_dict
    
    else:
        return render_template("create_pdf.html")


    
@app.route('/create_course', methods=["POST", "GET"])
def create_course():
    if request.method == "POST":
        course_name = request.form.get('course')
        num_quiz = request.form.get('quizzes')
        prompt = """generate %s questions in %s as json object with the keys question, answers, correct_answer 
it should be  in {} enclosed in a python list. I WANT ONLY THE OUTPUT""" % (num_quiz, course_name)
        print(prompt)
        prompt_results = generate_text(prompt)
        quiz_dict = eval(prompt_results)
        create_quiz_objects(quiz_dict, course_name)
        
        print(type(quiz_dict))
        return quiz_dict
    else:
        return render_template("create.html")



with app.app_context():
    db.create_all()

if __name__ == '__main__':
    app.run(debug=True)